"""AI 활용 고객경험관리시스템 활용 가이드 PPT 생성"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank

# ── 색상 정의 ──
NAVY = RGBColor(0x00, 0x2B, 0x5C)
BLUE = RGBColor(0x00, 0x55, 0xA5)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_BG = RGBColor(0xF0, 0xF4, 0xF8)
GOLD = RGBColor(0xFF, 0xB3, 0x00)
RED = RGBColor(0xE5, 0x3E, 0x30)
GREEN = RGBColor(0x2E, 0x7D, 0x32)
TEAL = RGBColor(0x00, 0x89, 0x7B)
PINK = RGBColor(0xE9, 0x1E, 0x63)
GRAY = RGBColor(0x66, 0x66, 0x66)
DARK = RGBColor(0x33, 0x33, 0x33)


def add_shape(slide, left, top, width, height, fill_color, border_color=None, radius=None):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape


def add_text(slide, left, top, width, height, text, font_size=12, color=DARK,
             bold=False, align=PP_ALIGN.LEFT, font_name="맑은 고딕"):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = align
    return txBox


def add_multiline(slide, left, top, width, height, lines, font_size=10, color=DARK,
                  bold=False, line_spacing=1.3, font_name="맑은 고딕"):
    """lines: list of (text, bold, color) or str"""
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        if isinstance(line, str):
            txt, b, c = line, bold, color
        else:
            txt, b, c = line
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = txt
        p.font.size = Pt(font_size)
        p.font.color.rgb = c
        p.font.bold = b
        p.font.name = font_name
        p.space_after = Pt(font_size * (line_spacing - 1) * 2)
    return txBox


# ═══════════════════════════════════════════════════════════
#  배경
# ═══════════════════════════════════════════════════════════
bg = slide.background
fill = bg.fill
fill.solid()
fill.fore_color.rgb = WHITE

# ── 상단 헤더 바 ──
header = add_shape(slide, 0, 0, 13.333, 1.05, NAVY)
add_text(slide, 0.5, 0.15, 8, 0.45, "AI 활용 고객경험관리시스템 분석 대시보드", 22, WHITE, True)
add_text(slide, 0.5, 0.58, 6, 0.35, "활용 가이드  |  경남본부 CS 리포트 분석 플랫폼", 11, RGBColor(0xBB, 0xCC, 0xDD))

# 우측 URL
add_text(slide, 8.5, 0.25, 4.5, 0.55,
         "huggingface.co/spaces/kepcokn5701/customer", 10, RGBColor(0xFF, 0xB3, 0x00),
         bold=True, align=PP_ALIGN.RIGHT)

# ═══════════════════════════════════════════════════════════
#  STEP 가이드 (상단)
# ═══════════════════════════════════════════════════════════
add_text(slide, 0.5, 1.2, 3, 0.35, "시작하기", 14, NAVY, True)

steps = [
    ("STEP 1", "엑셀 업로드", "좌측 사이드바에서\nCS 조사 엑셀 파일 업로드"),
    ("STEP 2", "컬럼 매핑", "지사·종합점수·업무구분\n·계약종별·VOC 컬럼 선택"),
    ("STEP 3", "분석 확인", "5개 탭에서\n자동 분석 결과 확인"),
]

for i, (step, title, desc) in enumerate(steps):
    x = 0.5 + i * 2.35
    box = add_shape(slide, x, 1.55, 2.15, 1.15, LIGHT_BG, RGBColor(0xDD, 0xDD, 0xDD))
    add_text(slide, x + 0.1, 1.58, 0.8, 0.25, step, 8, BLUE, True)
    add_text(slide, x + 0.1, 1.78, 1.95, 0.28, title, 12, NAVY, True)
    add_text(slide, x + 0.1, 2.05, 1.95, 0.6, desc, 9, GRAY)
    # 화살표
    if i < 2:
        add_text(slide, x + 2.15, 1.85, 0.25, 0.35, "→", 16, BLUE, True)

# ═══════════════════════════════════════════════════════════
#  5개 탭 기능 소개
# ═══════════════════════════════════════════════════════════
add_text(slide, 0.5, 2.9, 5, 0.35, "주요 기능 (5개 탭)", 14, NAVY, True)

tabs_data = [
    {
        "icon": "📊", "title": "종합 현황",
        "color": BLUE, "bg": RGBColor(0xE3, 0xF2, 0xFD),
        "items": [
            "종합만족도 게이지 (100점 환산)",
            "점수 구간별 비중 분석 (도넛 차트)",
            "핵심 KPI: 응답률·긍정비율·조기경보",
            "사업소별 평균 만족도 비교",
        ]
    },
    {
        "icon": "📡", "title": "계약·업무·항목별 분석",
        "color": TEAL, "bg": RGBColor(0xE0, 0xF2, 0xF1),
        "items": [
            "계약종별/업무유형별 만족도 비교",
            "지사×카테고리 교차 히트맵",
            "사업소별 항목별 결과표 (양식1)",
            "업무유형별 4분면 버블차트",
        ]
    },
    {
        "icon": "🏢", "title": "지사 맞춤형 CS 솔루션",
        "color": GREEN, "bg": RGBColor(0xE8, 0xF5, 0xE9),
        "items": [
            "KPI 헤더: 순위·피어그룹·신뢰도",
            "업무별 강점/약점 레이더 차트",
            "고객군 미스매치 버블 분석",
            "AI 리스크 분류 + 맞춤 솔루션",
        ]
    },
    {
        "icon": "🎯", "title": "민원 조기 경보",
        "color": RED, "bg": RGBColor(0xFF, 0xEB, 0xEE),
        "items": [
            "잠재 민원고객 자동 추출",
            "기준: 50점 이하 + 부정 키워드",
            "불만 유형별 도메인 키워드 분석",
            "사전케어 리스트 + AI 리스크 진단",
        ]
    },
    {
        "icon": "💌", "title": "경험고객 서한문 생성",
        "color": PINK, "bg": RGBColor(0xFC, 0xE4, 0xEC),
        "items": [
            "지사·계절 선택 → AI 맞춤 서한문",
            "VOC 반영 톤 자동 조절 (Gemini Pro)",
            "서한문 디자인 미리보기 + 문구 복사",
            "계절·지역·VOC별 기념품 3종 추천",
        ]
    },
]

card_w = 2.35
card_h = 2.55
gap = 0.15
start_x = 0.5
start_y = 3.3

for i, tab in enumerate(tabs_data):
    x = start_x + i * (card_w + gap)

    # 카드 배경
    add_shape(slide, x, start_y, card_w, card_h, tab["bg"], tab["color"])

    # 아이콘 + 제목
    add_text(slide, x + 0.12, start_y + 0.08, card_w - 0.2, 0.35,
             f'{tab["icon"]}  {tab["title"]}', 12, tab["color"], True)

    # 구분선
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(x + 0.12), Inches(start_y + 0.45),
        Inches(card_w - 0.24), Pt(1.5))
    line.fill.solid()
    line.fill.fore_color.rgb = tab["color"]
    line.line.fill.background()

    # 항목들
    items_lines = [(f"•  {item}", False, DARK) for item in tab["items"]]
    add_multiline(slide, x + 0.12, start_y + 0.55, card_w - 0.24, card_h - 0.7,
                  items_lines, font_size=9, line_spacing=1.5)

# ═══════════════════════════════════════════════════════════
#  하단 안내
# ═══════════════════════════════════════════════════════════
# 좌측: AI 모델 정보
add_shape(slide, 0.5, 6.15, 6.0, 1.1, RGBColor(0xF8, 0xF9, 0xFA), RGBColor(0xDD, 0xDD, 0xDD))
add_text(slide, 0.7, 6.2, 3, 0.25, "AI 엔진", 10, NAVY, True)
add_multiline(slide, 0.7, 6.45, 5.6, 0.7, [
    ("•  분석/솔루션:  Gemini 2.5 Flash  →  Gemma 3 폴백", False, GRAY),
    ("•  서한문 생성:  Gemini 2.5 Pro  →  Flash-Lite 폴백", False, GRAY),
    ("•  감성분석: 룰 기반 한국어 키워드 분류 (외부 API 불필요)", False, GRAY),
], font_size=9, line_spacing=1.4)

# 우측: 참고사항
add_shape(slide, 6.7, 6.15, 6.15, 1.1, RGBColor(0xF8, 0xF9, 0xFA), RGBColor(0xDD, 0xDD, 0xDD))
add_text(slide, 6.9, 6.2, 3, 0.25, "참고사항", 10, NAVY, True)
add_multiline(slide, 6.9, 6.45, 5.8, 0.7, [
    ("•  엑셀 필수 컬럼: 지사, 종합점수, 서술의견(VOC) — 나머지는 선택", False, GRAY),
    ("•  VOC 컬럼 선택 시 감성분석·조기경보·AI 솔루션 활성화", False, GRAY),
    ("•  서한문은 인터넷 연결 필요 (Gemini API 호출)", False, GRAY),
], font_size=9, line_spacing=1.4)


# ═══════════════════════════════════════════════════════════
#  저장
# ═══════════════════════════════════════════════════════════
output = r"c:\Users\Admin\Desktop\customer\활용가이드_AI고객경험관리시스템.pptx"
prs.save(output)
print(f"PPT 생성 완료: {output}")
